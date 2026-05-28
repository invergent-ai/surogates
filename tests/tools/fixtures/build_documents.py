"""Programmatic builders for tiny office documents used in tests.

Generates fixtures at test time rather than checking binaries into git.
Each builder writes a small but real document with a known probe string
so tests can assert that liteparse extracted it.
"""

from __future__ import annotations

from pathlib import Path


def build_minimal_pdf(path: Path, heading: str = "Hello PDF") -> Path:
    """Write a minimal 2-page PDF with ``heading`` on each page.

    Uses ``reportlab`` so the fixture contains a real extractable text
    stream (a raw pypdf-only PDF with no text wouldn't round-trip).
    """
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(path))
    c.setFont("Helvetica", 14)
    c.drawString(72, 720, heading)
    c.showPage()
    c.setFont("Helvetica", 14)
    c.drawString(72, 720, f"{heading} page 2")
    c.showPage()
    c.save()
    return path


def build_minimal_docx(path: Path) -> Path:
    """Write a docx with one H1, one bullet, one 2x2 table."""
    from docx import Document

    doc = Document()
    doc.add_heading("Hello DOCX", level=1)
    doc.add_paragraph("first bullet", style="List Bullet")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"
    doc.save(str(path))
    return path


def build_minimal_xlsx(path: Path) -> Path:
    """Write an xlsx with two sheets, each containing a probe string in cells.

    Probes are written into cells (not sheet titles) because liteparse
    routes xlsx → LibreOffice → PDF, and the PDF stream doesn't carry
    sheet metadata.
    """
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Alpha"
    ws1["A1"] = "Alpha header"
    ws1["B1"] = "value"
    ws1["A2"] = "x"
    ws1["B2"] = 1
    ws2 = wb.create_sheet("Beta")
    ws2["A1"] = "Beta header"
    ws2["A2"] = "y"
    wb.save(str(path))
    return path


def build_minimal_pptx(path: Path) -> Path:
    """Write a pptx with one slide titled 'Hello PPTX'."""
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Hello PPTX"
    slide.placeholders[1].text = "body text"
    prs.save(str(path))
    return path
