"""Best-effort text extraction from collected output files for the judge.

The upstream judge is agentic -- it inspects the filesystem itself. Ours
is a single completion, so the file contents must be inlined, which means
office formats need extraction. Every extractor is fenced: an unreadable
file degrades to a note the judge can see ("binary, not extractable"),
never to a crashed grading run. That asymmetry is deliberate -- a rubric
judged against "the file exists but could not be read" fails with a
visible reason, which is the honest outcome.
"""
from __future__ import annotations

import posixpath

_CODE_TEXT_EXTS = {
    ".md", ".txt", ".csv", ".tsv", ".json", ".jsonl", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".xml", ".html", ".css", ".svg", ".sql",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".kt", ".c", ".cpp",
    ".h", ".hpp", ".go", ".rs", ".rb", ".php", ".sh", ".bash", ".r",
    ".ipynb", ".tex", ".rst", ".log", ".env", ".properties",
}

DEFAULT_MAX_CHARS = 30_000


def _truncate(text: str, max_chars: int) -> tuple[str, str | None]:
    if len(text) <= max_chars:
        return text, None
    return (
        text[:max_chars],
        f"truncated to {max_chars} of {len(text)} chars",
    )


def _extract_docx(path: str) -> str:
    import docx

    doc = docx.Document(path)
    parts = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _chart_title(chart) -> str:
    """Dig the plain text out of openpyxl's Title object graph."""
    title = getattr(chart, "title", None)
    if title is None:
        return "(untitled)"
    try:
        runs = []
        for p in title.tx.rich.p:
            for r in p.r or []:
                if r.t:
                    runs.append(r.t)
        return "".join(runs) or "(untitled)"
    except Exception:  # noqa: BLE001 - title shapes vary wildly
        return str(title)


def _extract_xlsx(path: str, max_rows_per_sheet: int = 300) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for ws in wb.worksheets:
            parts.append(f"# sheet: {ws.title}")
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= max_rows_per_sheet:
                    parts.append(f"... ({ws.max_row} rows total)")
                    break
                parts.append(
                    "\t".join("" if v is None else str(v) for v in row)
                )
    finally:
        wb.close()

    # Charts are invisible in read_only mode, and rubrics routinely ask
    # about chart title/type/placement -- without this, every chart
    # rubric fails as "no evidence" even when the agent did the work.
    try:
        wb2 = openpyxl.load_workbook(path)
        try:
            for ws in wb2.worksheets:
                for chart in getattr(ws, "_charts", []):
                    kind = type(chart).__name__
                    direction = getattr(chart, "type", None)
                    if kind == "BarChart":
                        # openpyxl: type "bar" = horizontal, "col" = vertical.
                        kind += f" ({'horizontal' if direction == 'bar' else 'vertical'})"
                    anchor = getattr(chart, "anchor", None)
                    cell = getattr(anchor, "_from", None)
                    where = (
                        f"row {cell.row + 1}, col {cell.col + 1}"
                        if cell is not None else str(anchor or "?")
                    )
                    parts.append(
                        f"# chart on sheet '{ws.title}': {kind}, "
                        f"title: {_chart_title(chart)}, "
                        f"series: {len(getattr(chart, 'series', []) or [])}, "
                        f"anchored at {where}"
                    )
        finally:
            wb2.close()
    except Exception as exc:  # noqa: BLE001 - charts are best-effort
        parts.append(f"# chart inspection failed ({type(exc).__name__})")

    return "\n".join(parts)


def _extract_pptx(path: str) -> str:
    import pptx

    prs = pptx.Presentation(path)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        parts.append(text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append(f"[notes] {notes}")
    return "\n".join(parts)


def _extract_pdf(path: str, max_pages: int = 30) -> str:
    import pypdf

    reader = pypdf.PdfReader(path)
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        if i >= max_pages:
            parts.append(f"... ({len(reader.pages)} pages total)")
            break
        parts.append(page.extract_text() or "")
    return "\n".join(parts)


def extract_text(
    local_path: str, max_chars: int = DEFAULT_MAX_CHARS
) -> tuple[str, str | None]:
    """Return (text, note). ``note`` explains truncation or failure."""
    ext = posixpath.splitext(local_path.lower())[1]
    try:
        if ext == ".docx":
            return _truncate(_extract_docx(local_path), max_chars)
        if ext == ".xlsx":
            return _truncate(_extract_xlsx(local_path), max_chars)
        if ext == ".pptx":
            return _truncate(_extract_pptx(local_path), max_chars)
        if ext == ".pdf":
            return _truncate(_extract_pdf(local_path), max_chars)
        if ext in _CODE_TEXT_EXTS or not ext:
            with open(local_path, encoding="utf-8", errors="replace") as fh:
                return _truncate(fh.read(), max_chars)
    except Exception as exc:  # noqa: BLE001 - degrade, never crash grading
        return "", f"extraction failed ({type(exc).__name__}: {exc})"

    # Unknown binary format: try utf-8, else report as binary.
    try:
        with open(local_path, "rb") as fh:
            blob = fh.read(max_chars * 4)
        text = blob.decode("utf-8")
        return _truncate(text, max_chars)
    except UnicodeDecodeError:
        return "", f"binary file ({ext or 'no extension'}), content not extractable"
